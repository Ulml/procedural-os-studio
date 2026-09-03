"""Generate md / docx / xlsx / pptx / charts / placeholder images. No Google APIs."""

from __future__ import annotations

from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from .paths import ARTIFACTS


def _safe(name: str, suffix: str) -> Path:
    stem = "".join(ch if ch.isalnum() or ch in "-_" else "-" for ch in name).strip("-") or "artifact"
    path = ARTIFACTS / f"{stem}{suffix}"
    ARTIFACTS.mkdir(parents=True, exist_ok=True)
    return path


def write_markdown(name: str, title: str, body: str) -> Path:
    path = _safe(name, ".md")
    path.write_text(f"# {title}\n\n{body.strip()}\n", encoding="utf-8")
    return path


def write_docx(name: str, title: str, body: str) -> Path:
    path = _safe(name, ".docx")
    from docx import Document

    doc = Document()
    doc.add_heading(title, 0)
    for para in body.split("\n\n"):
        doc.add_paragraph(para.strip())
    doc.add_paragraph(f"Generated {datetime.now(timezone.utc).isoformat()}")
    doc.save(str(path))
    return path


def write_xlsx(name: str, title: str, rows: list[list[Any]] | None = None) -> Path:
    path = _safe(name, ".xlsx")
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = title[:31] or "Sheet"
    for r in rows or [["metric", "value"], [title, 1]]:
        ws.append(list(r))
    wb.save(str(path))
    return path


def write_pptx(name: str, title: str, bullets: list[str] | None = None) -> Path:
    path = _safe(name, ".pptx")
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, line in enumerate(bullets or [title]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18)
    prs.save(str(path))
    return path


def write_chart(name: str, labels: list[str], values: list[float], title: str = "Chart") -> Path:
    path = _safe(name, ".png")
    import matplotlib

    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    fig, ax = plt.subplots(figsize=(8, 4.5))
    ax.bar(labels, values, color="#C6A46A")
    ax.set_title(title)
    fig.tight_layout()
    fig.savefig(path, dpi=140)
    plt.close(fig)
    return path


def write_image(name: str, title: str) -> Path:
    path = _safe(name, ".png")
    from PIL import Image, ImageDraw

    img = Image.new("RGB", (1280, 720), "#1E2A3A")
    draw = ImageDraw.Draw(img)
    draw.rectangle([40, 40, 1240, 680], outline="#C6A46A", width=3)
    draw.text((80, 300), title[:80], fill="#EFE9DD")
    img.save(path)
    return path


def build(kind: str, name: str, payload: dict) -> dict:
    title = payload.get("title") or name
    body = payload.get("body") or ""
    if kind == "md":
        path = write_markdown(name, title, body)
    elif kind == "docx":
        path = write_docx(name, title, body)
    elif kind == "xlsx":
        path = write_xlsx(name, title, payload.get("rows"))
    elif kind == "pptx":
        path = write_pptx(name, title, payload.get("bullets") or body.split("\n"))
    elif kind == "chart":
        path = write_chart(name, payload.get("labels") or ["A", "B"], payload.get("values") or [1, 2], title)
    elif kind == "image":
        path = write_image(name, title)
    else:
        raise ValueError(f"unknown artifact kind {kind}")
    return {"ok": True, "kind": kind, "path": str(path), "name": path.name}
